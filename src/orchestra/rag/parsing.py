"""文档解析：从 Markdown/文本/PDF/Word/Excel/PPT 中提取可索引文本。

解析器保持轻量依赖：只有对应文件类型存在时才导入对应库，
便于在未安装全部解析依赖时仍可处理纯文本/Markdown。
"""
from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

_SUPPORTED_SUFFIXES = {
    ".md",
    ".markdown",
    ".txt",
    ".text",
    ".pdf",
    ".docx",
    ".xlsx",
    ".xlsm",
    ".pptx",
}


@dataclass(frozen=True)
class ParsedSection:
    """一份文档的一节文本，page 用于 PDF 分页溯源。"""

    text: str
    page: int | None = None


def _read_text(path: Path) -> str:
    """按 UTF-8 / GB18030 顺序读取文本文件，避免中文编码不一致。"""
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _clean(text: str) -> str:
    """合并连续空白，保留换行，便于后续分块与检索。"""
    return re.sub(r"[ \t\u3000]+", " ", text.replace("\r\n", "\n").replace("\r", "\n"))


def _parse_pdf(path: Path) -> list[ParsedSection]:
    """按页提取 PDF 文本，保留页码元数据。"""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    sections: list[ParsedSection] = []
    for page_index, page in enumerate(reader.pages, start=1):
        text = _clean(page.extract_text() or "")
        if text.strip():
            sections.append(ParsedSection(text=text, page=page_index))
    return sections


def _parse_docx(path: Path) -> list[ParsedSection]:
    """解析 Word：段落 + 表格，表格行用竖线拼接。"""
    from docx import Document

    document = Document(str(path))
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return [ParsedSection(text="\n".join(parts))] if parts else []


def _parse_xlsx(path: Path) -> list[ParsedSection]:
    """解析 Excel：每个工作表输出为一节，空行跳过。"""
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    sections: list[ParsedSection] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value).strip() for value in row]
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            sections.append(ParsedSection(text=f"工作表：{sheet.title}\n" + "\n".join(rows)))
    workbook.close()
    return sections


def _parse_pptx(path: Path) -> list[ParsedSection]:
    """解析 PPT：按幻灯片提取文本与表格。"""
    from pptx import Presentation

    presentation = Presentation(str(path))
    sections: list[ParsedSection] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if parts:
            sections.append(ParsedSection(text="\n".join(parts), page=slide_index))
    return sections


def parse_document(path: Path) -> list[ParsedSection]:
    """按扩展名解析单个文档，返回可索引的文本节列表。"""
    suffix = path.suffix.lower()
    if suffix not in _SUPPORTED_SUFFIXES:
        raise ValueError(f"不支持的文件类型：{suffix}，仅支持 {sorted(_SUPPORTED_SUFFIXES)}")
    if suffix in {".md", ".markdown", ".txt", ".text"}:
        return [ParsedSection(text=_clean(_read_text(path)))]
    if suffix == ".pdf":
        return _parse_pdf(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix in {".xlsx", ".xlsm"}:
        return _parse_xlsx(path)
    if suffix == ".pptx":
        return _parse_pptx(path)
    return []
